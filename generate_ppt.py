import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Antique Marketplace System"
subtitle.text = "Database Design, ER Diagram, and Flowchart\nCreated via AI"

# Helper to add a slide with image
def add_image_slide(prs, title_text, image_path):
    blank_slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = title_text
    
    # Calculate image position to center it
    try:
        left = Inches(1)
        top = Inches(2)
        height = Inches(4.5)
        # Using height ensures aspect ratio is maintained
        pic = slide.shapes.add_picture(image_path, left, top, height=height)
    except Exception as e:
        print(f"Error adding {image_path}: {e}")

# Slide 1: Database Design
add_image_slide(prs, "Database Design", r"C:\Users\wel\.gemini\antigravity\brain\a0279335-49ed-4e49-8a49-5bdb5c801a56\antique_database_design_1777522841995.png")

# Slide 2: ER Diagram
add_image_slide(prs, "Entity-Relationship (ER) Diagram", r"C:\Users\wel\.gemini\antigravity\brain\a0279335-49ed-4e49-8a49-5bdb5c801a56\antique_er_diagram_1777522827306.png")

# Slide 3: Application Flowchart
add_image_slide(prs, "Application Flowchart", r"C:\Users\wel\.gemini\antigravity\brain\a0279335-49ed-4e49-8a49-5bdb5c801a56\antique_flowchart_1777522857582.png")

# Save Presentation
output_path = r"C:\Users\wel\.gemini\antigravity\brain\a0279335-49ed-4e49-8a49-5bdb5c801a56\artifacts\Antique_Marketplace_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
